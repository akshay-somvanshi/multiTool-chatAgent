from langchain_core.tools import tool, Tool, StructuredTool
from google import genai
from fpdf import FPDF
from langchain_google_community import GoogleSearchAPIWrapper
from pydantic import BaseModel, Field
from google.api_core.client_options import ClientOptions
from google.cloud import discoveryengine_v1 as discoveryengine
from google.cloud import documentai_v1
import os
import dotenv
import requests
import google
from urllib.parse import urlparse, parse_qs, unquote
from google.cloud import storage
from google.cloud import bigquery
from datetime import datetime
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
import pandas as pd
import io
import json
import contextvars

# Per-request Firestore context — set by the agent at the start of each request.
# ContextVar is safe for concurrent async tasks: each task (and threads it spawns)
# gets an independent copy, so simultaneous users don't interfere.
_firestore_ctx: contextvars.ContextVar = contextvars.ContextVar(
    'firestore_ctx', default=None
)

from chat_agent.api_client import api_client
from chat_agent.core.exceptions import APIError
from chat_agent.core.octopus import OctopusClient
from chat_agent.firestore import FireStoreChat

dotenv.load_dotenv()

class google_search_input(BaseModel):
    search_query: str = Field(
        description='The natural language query to search on google'
    )

class vertex_search_input(BaseModel):
    query: str = Field(
        description='The natural language query to search against the data store'
    ),
    user_id: str = Field(
        description='The id of the user whos documents can be accessed. '
    )

class document_read_input(BaseModel):
    document_url: str = Field(
        description='The url of the document that must be read'
    )

class search_query(BaseModel):
    search_query: str = Field(
        description='The natural language query to process'
    )

class get_api(BaseModel):
    user_id: str = Field(
        description="User id whos actions are fetched from the database"
    )

class EnergyFetchInput(BaseModel):
    user_id: str = Field(description="The user ID")
    days_back: int = Field(default=7, description="Number of days to fetch consumption for (used if period_from is not set)")
    period_from: str | None = Field(default=None, description="Start of the period in ISO format (e.g. 2026-01-01T00:00:00Z)")
    period_to: str | None = Field(default=None, description="End of the period in ISO format (e.g. 2026-02-01T00:00:00Z)")

class AddActionInput(BaseModel):
    user_id: str = Field(description="User ID")
    action_id: str = Field(description="Unique action identifier of type 'action_003'")
    action_name: str = Field(description="Name of the action")
    action_type: str = Field(description="Type of action (e.g., energy_efficiency)")
    action_description: str = Field(description="Detailed description")
    estimated_spend: float = Field(description="Estimated cost in GBP")
    estimated_co2_reduced: float = Field(description="Estimated CO2 reduction in tCO2e/year")
    estimated_revenue_unlocked: float = Field(description="Estimated Revenue potential in GBP")
    plan_id: str = Field(description="Associated plan ID")
    timeline_start: datetime = Field(description="Proposed start date")
    timeline_end: datetime = Field(description="Proposed end date")
    status: str = Field(description="Current status. Default : not_started")

class UpdateActionInput(BaseModel):
    user_id: str = Field(description="User ID")
    action_id: str = Field(description="Unique action identifier of type 'action_003'")
    actual_co2_reduced: float | None = Field(description="Actual CO2 reduced in tCO2e")
    actual_spend: float | None = Field(description="Actual spend in GBP")
    actual_revenue_unlocked: float | None = Field(description="Actual revenue unlocked in GBP")
    day_started: datetime | None = Field(description="Date when action was started")
    day_completed: datetime | None = Field(description="Date when action was completed")

class RemoveActionInput(BaseModel):
    user_id: str = Field(description="User ID")
    action_id: str = Field(description="Unique action identifier of type 'action_003'")

class ProcurementROIInput(BaseModel):
    user_id: str = Field(description="User ID")
    action_description: str = Field(description="Description of the sustainability action being evaluated for procurement eligibility")

class PolicyGapsInput(BaseModel):
    user_id: str = Field(description="User ID to check policy gaps for")

class DecarbonisationActionsInput(BaseModel):
    user_id: str = Field(description="User ID to fetch decarbonisation actions for")

class IndustryInfoInput(BaseModel):
    industry_name: str = Field(description="The name of the industry to look up guidelines for (e.g., 'Retail', 'Manufacturing')")

class PDFGeneratorInput(BaseModel):
    title: str = Field(description="The title shown inside the report")
    content: str = Field(description="The text content of the report. Use \n for new lines.")
    user_id: str = Field(description="The user ID to associate the report with")
    company_name: str = Field(description="The name of the company (used for filename)")
    report_type: str = Field(description="The type of report, e.g. 'Sustainability plan', 'CDP document' (used for filename)")

class PPTXSlide(BaseModel):
    title: str = Field(description="Slide heading (short, max 8 words)")
    bullets: list[str] = Field(description="3-6 bullet points for the slide body. Each bullet is one concise sentence.")
    notes: str = Field(default="", description="Optional speaker notes for this slide")

class PPTXGeneratorInput(BaseModel):
    title: str = Field(description="Presentation title (shown on the cover slide)")
    subtitle: str = Field(default="", description="Optional subtitle or tagline for the cover slide")
    slides: list[PPTXSlide] = Field(description="List of content slides. Each slide has a title and bullets. Provide 4-10 slides.")
    user_id: str = Field(description="The user ID to associate the file with")
    company_name: str = Field(description="Company name — shown on the cover and used in the filename")
    report_type: str = Field(description="Type label, e.g. 'Sustainability Plan', 'Emissions Overview' — used in filename")

class CarbonCalculationInput(BaseModel):
    activity_name: str = Field(description="The name of the activity (e.g., 'Electricity', 'Petrol', 'Natural Gas', 'Short-haul Flight')")
    amount: float = Field(description="The numerical value of the activity (e.g., 500)")
    unit: str = Field(description="The unit for the activity (e.g., 'kWh', 'litres', 'passenger-km', 'tonnes')")

class BulkReadinessInput(BaseModel):
    user_id: str = Field(description="The user ID whose uploads folder should be checked")
    expected_categories: list[str] = Field(default=[], description="Optional list of categories the user expects to find (e.g. ['Electricity', 'Fuel'])")

class BulkProcessInput(BaseModel):
    user_id: str = Field(description="The user ID whose uploaded file should be processed")

class ToolList:
    def __init__(self):
        self.project_id = os.getenv("GOOGLE_PROJECT_ID")
        self.location = os.getenv("GOOGLE_LOCATION")
        self.location_vertexAI = "eu"
        self.engine_id = os.getenv("VERTEX_ENGINE_ID")
        self.document_ai_id = os.getenv("DOCUMENT_AI_ID")

        # Initialise document AI client
        # Set api endpoint to eu
        opts = ClientOptions(api_endpoint=f"{self.location_vertexAI}-documentai.googleapis.com")
        self.docai_client = documentai_v1.DocumentProcessorServiceClient(client_options=opts)
        # Processor reference
        self.processor_name = self.docai_client.processor_path(
            self.project_id, 
            self.location_vertexAI,
            self.document_ai_id
        )

        # Initialise discovery engine client
        client_options = (
            ClientOptions(api_endpoint=f"{self.location_vertexAI}-discoveryengine.googleapis.com")
            if self.location_vertexAI != "global"
            else None
        )
        self.discovery_engine_client = discoveryengine.ConversationalSearchServiceClient(client_options=client_options) 
        self.bq_client = bigquery.Client(project=self.project_id)
        self.genai_client = genai.Client(
            vertexai=True,
            project=self.project_id,
            location="global"
        )

    def _emit_status(self, key: str) -> None:
        """Write a live status update to Firestore via the per-request context."""
        fs = _firestore_ctx.get()
        if fs is not None:
            try:
                fs.set_status(key)
            except Exception as e:
                print(f"[Status] Failed to emit status '{key}': {e}", flush=True)

    def _logged_search(
        self,
        query
    ):
        self._emit_status("google_search")
        print(f"Google called with query: {query}")
        # Create a fresh wrapper to avoid httplib2 SSL socket reuse bugs on Cloud Run
        fresh_wrapper = GoogleSearchAPIWrapper()
        result = fresh_wrapper.run(query)
        print(f"Search result: {result[:200]}...")
        return result

    def _download_from_gcs(self, url):
        """
        Download a file from Google Cloud Storage using the standard storage client.
        
        Supports multiple URL formats:
        - Authenticated Console URLs (storage.cloud.google.com/bucket/blob)
        - Public API URLs (storage.googleapis.com/bucket/blob)
        - Subdomain API URLs (bucket.storage.googleapis.com/blob)
        - Firebase Storage URLs (bucket.firebasestorage.app/blob)
        - Native gs:// URI format (gs://bucket/blob)

        If the file has a .txt or .csv extension, it is returned as a UTF-8 string. 
        Otherwise, it is returned as raw bytes.

        Args:
            url (str): The GCS URL or gs:// URI to download.

        Returns:
            Union[str, bytes]: File content as text (for .txt/.csv) or bytes (other formats).
        """
        try:
            # Handle native gs:// format
            if url.startswith("gs://"):
                parts = url[5:].split('/', 1)
                bucket_name = parts[0]
                blob_path = parts[1] if len(parts) > 1 else ""
            else:
                parsed = urlparse(url)
                netloc = parsed.netloc
                path = parsed.path.lstrip('/')
                
                # Extract bucket and blob based on host
                if 'storage.cloud.google.com' in netloc:
                    bucket_name, blob_path = path.split('/', 1)
                elif 'storage.googleapis.com' in netloc:
                    if netloc == 'storage.googleapis.com':
                        bucket_name, blob_path = path.split('/', 1)
                    else:
                        bucket_name = netloc.split('.')[0]
                        blob_path = path
                elif 'firebasestorage.app' in netloc:
                    bucket_name = netloc.split('.')[0]
                    blob_path = path
                else:
                    raise ValueError(f"Unsupported GCS URL format: {url}")

            # Decode URL-encoded characters in the blob path (like %20 for spaces)
            blob_path = unquote(blob_path)
            print(f"Downloading from GCS: bucket={bucket_name}, path={blob_path}", flush=True)
            
            storage_client = storage.Client(project=self.project_id)
            bucket = storage_client.bucket(bucket_name)
            blob = bucket.blob(blob_path)
            
            # Simple text vs bytes logic
            # .xlsx is a binary format, MUST be bytes
            if blob_path.lower().endswith('.txt') or blob_path.lower().endswith('.csv'):
                print(f"Downloading {blob_path} as text", flush=True)
                content = blob.download_as_text()
            else:
                print(f"Downloading {blob_path} as bytes", flush=True)
                content = blob.download_as_bytes()
            
            return content
            
        except Exception as e:
            if url.startswith("gs://"):
                raise  # gs:// URIs cannot be fetched over HTTP
            print(f"GCS download failed: {e}. Falling back to HTTP.", flush=True)
            return self._download_from_http(url)

    def _download_from_http(self, url):
        """Download file via HTTP."""
        headers = {
            'User-Agent': 'Mozilla/5.0 (compatible; DocumentAI/1.0)'
        }
        
        response = requests.get(
            url, 
            timeout=60,
            headers=headers,
            stream=True
        )
        response.raise_for_status()
        return response.content
    
    def _excel_to_text(self, content: bytes) -> str:
        """
        Reads all sheets from an Excel file and returns their combined text.
        Each sheet is converted to CSV in parallel, then joined with a header.
        """
        sheets = pd.read_excel(io.BytesIO(content), sheet_name=None)
        non_empty = {name: df for name, df in sheets.items() if not df.empty}
        if not non_empty:
            return "(Excel file contained no data)"

        print(f"[Tool] Converting {len(non_empty)} sheet(s) to text in parallel", flush=True)

        def sheet_to_csv(item):
            name, df = item
            return f"=== Sheet: {name} ===\n{df.to_csv(index=False)}"

        with ThreadPoolExecutor(max_workers=min(len(non_empty), 5)) as executor:
            results = list(executor.map(sheet_to_csv, non_empty.items()))

        return "\n\n".join(results)

    def _document_read(
        self,
        document_url
    ):
        self._emit_status("document_read")
        print(f"[Tool] Starting document_read for: {document_url}", flush=True)
        """
        Reads and extracts text from a document (PDF, TXT, CSV).
        
        - If the document is a .txt or .csv it is read directly.
        - If the document is a PDF, it is processed via Google Document AI OCR.
        - Supports GCS URLs (authenticated, public, firebase) and HTTP links.

        Args:
            document_url (str): URL or path to the document.

        Returns:
            str: Extracted text content.
        """
        try:
            gcs_identifiers = ['storage.googleapis.com', 'firebasestorage.app', 'storage.cloud.google.com', 'gs://']
            is_gcs = any(id in document_url for id in gcs_identifiers)

            if is_gcs:
                print(f"Detected GCS document", flush=True)
                # Skip Document AI for native text formats
                if any(ext in document_url.lower() for ext in ['.txt', '.csv']):
                    print(f"Skipping Document AI for text-based file", flush=True)
                    return self._download_from_gcs(document_url)
                elif '.xlsx' in document_url.lower():
                    print(f"Reading Excel file (all sheets) for text extraction", flush=True)
                    content = self._download_from_gcs(document_url)
                    return self._excel_to_text(content)
                else:
                    image_content = self._download_from_gcs(document_url)
            elif document_url.startswith("http://") or document_url.startswith("https://"):
                print(f"Downloading file via HTTP", flush=True)
                # Handle direct Excel/CSV over HTTP as well
                if any(ext in document_url.lower() for ext in ['.txt', '.csv']):
                    return self._download_from_http(document_url).decode('utf-8', errors='replace')
                elif '.xlsx' in document_url.lower():
                    content = self._download_from_http(document_url)
                    return self._excel_to_text(content)
                
                image_content = self._download_from_http(document_url)
            else:
                if not os.path.exists(document_url):
                    raise FileNotFoundError(f"File not found: {document_url}")
                
                with open(document_url, "rb") as f:
                    image_content = f.read()
                print(f"Local file loaded (size: {len(image_content) / 1_024:.2f} KB)", flush=True)

            raw_doc = documentai_v1.RawDocument(
                content=image_content,
                mime_type="application/pdf"
            )

            # Verify processor name
            print(f"Using processor: {self.processor_name}", flush=True)

            # Create request
            request = documentai_v1.ProcessRequest(
                name=self.processor_name, 
                raw_document=raw_doc
            )
            print("ProcessRequest created, calling Document AI...", flush=True)
            
            # Call Document AI with explicit error handling
            try:
                result = self.docai_client.process_document(
                    request=request,
                    timeout=600.0
                )
                print("Document AI call completed successfully", flush=True)
                
            except google.api_core.exceptions.GoogleAPICallError as e:
                print(f"Document AI API error: {e}", flush=True)
                print(f"Error details: {e.details()}", flush=True)
                raise
            except google.api_core.exceptions.NotFound as e:
                print(f"Processor not found: {e}", flush=True)
                raise
            except google.api_core.exceptions.PermissionDenied as e:
                print(f"Permission denied: {e}", flush=True)
                raise
            except google.api_core.exceptions.DeadlineExceeded as e:
                print(f"Request deadline exceeded: {e}", flush=True)
                raise
            except Exception as e:
                print(f"Unexpected error during Document AI call: {type(e).__name__}", flush=True)
                print(f"Error message: {str(e)}", flush=True)
                import traceback
                print(f"Traceback: {traceback.format_exc()}", flush=True)
                raise
            
            document = result.document
            print(f"Text extraction complete, length: {len(document.text)} characters", flush=True)
            
            return document.text
        
        except Exception as e:
            print(f"Fatal error in _document_read: {type(e).__name__}: {str(e)}", flush=True)
            import traceback
            print(f"Full traceback: {traceback.format_exc()}", flush=True)
            raise
    
    def read_actions(
            self,
            user_id: str
    ) -> dict:
        """ Fetches all sustainability actions for the user. Use this when you want to look up the actions in the database.W"""
        try:
            return api_client.view_actionList(user_id)
        except Exception as e:
            return {
                "error": {e},
                "actions": []
            }

    def calculate_procurement_roi(self, user_id: str, action_description: str) -> dict:
        """
        Calculates ROI as a procurement eligibility score: how many industry procurement
        policies does this action satisfy? Policies are stored in BigQuery and sourced from
        publicly available industry procurement standards.
        """
        try:
            # Fetch user's industry from Firestore
            session_id = f"roi_{datetime.now().strftime('%Y%m%d')}"
            fs = FireStoreChat(user_id, session_id)
            user_doc = fs.user_ref.get()
            if not user_doc.exists:
                return {"error": "User not found"}

            user_data = user_doc.to_dict()
            industry = user_data.get("company_industry", "").strip()
            if not industry:
                return {"error": "User's industry is not set in their profile"}

            # Fetch procurement policies for this industry from BigQuery
            dataset_id = "dash_beta_database"
            table_id = "procurement_policies"
            query = f"""
                SELECT policy_id, policy_name, policy_description, category, policy_source,
                       IFNULL(policy_type, 'regulatory') as policy_type,
                       source_company
                FROM `{self.project_id}.{dataset_id}.{table_id}`
                WHERE LOWER(industry) = @industry
                ORDER BY policy_type, category, policy_name
            """
            job_config = bigquery.QueryJobConfig(
                query_parameters=[
                    bigquery.ScalarQueryParameter("industry", "STRING", industry.lower())
                ]
            )
            results = list(self.bq_client.query(query, job_config=job_config).result())

            if not results:
                return {
                    "error": f"No procurement policies found for industry: {industry}",
                    "hint": "Run scripts/populate_procurement_policies.py to seed policies for this industry."
                }

            policies = [
                {
                    "policy_id": row.policy_id,
                    "policy_name": row.policy_name,
                    "policy_description": row.policy_description,
                    "category": row.category,
                    "policy_source": row.policy_source,
                    "policy_type": row.policy_type,
                    "source_company": row.source_company or None,
                }
                for row in results
            ]
            # Build a lookup so we can enrich evaluations with policy_source on the way out
            policy_lookup = {p["policy_id"]: p for p in policies}

            prompt = f"""You are evaluating a sustainability action against procurement policies for the {industry} industry.

Action being evaluated:
{action_description}

Procurement policies to assess (mix of regulatory requirements and corporate standards from leading industry companies):
{json.dumps(policies, indent=2)}

For each policy, determine whether completing this action would directly and fully satisfy that policy requirement.
- For regulatory policies: mark satisfied only if the action clearly meets the legal/standard requirement.
- For corporate policies: mark satisfied if the action meets what that specific company requires of its suppliers.
Mark 'satisfied' as true only if the action clearly and completely meets the criterion — partial alignment does not count.

You MUST return an evaluation for every policy in the list above. Do not skip any.

Return ONLY a valid JSON object in this exact structure:
{{
  "evaluations": [
    {{
      "policy_id": "<id>",
      "policy_name": "<name>",
      "category": "<category>",
      "satisfied": true or false,
      "reason": "<one sentence explanation>"
    }}
  ]
}}"""

            evaluations = []
            for attempt in range(3):
                try:
                    response = self.genai_client.models.generate_content(
                        model="gemini-3.5-flash",
                        contents=prompt,
                        config={"response_mime_type": "application/json"}
                    )
                    evaluations = json.loads(response.text).get("evaluations", [])
                    break
                except Exception as llm_err:
                    print(f"[ROI] LLM eval attempt {attempt + 1} failed: {llm_err}", flush=True)
                    if attempt == 2:
                        return {"error": f"LLM evaluation failed: {llm_err}"}
                    time.sleep(2 ** attempt)

            # Use the BQ policy count as the authoritative denominator.
            # Any policy the LLM omitted is treated as unsatisfied.
            total = len(policies)
            evaluated_ids = {e.get("policy_id") for e in evaluations}
            for p in policies:
                if p["policy_id"] not in evaluated_ids:
                    evaluations.append({
                        "policy_id": p["policy_id"],
                        "policy_name": p["policy_name"],
                        "category": p["category"],
                        "satisfied": False,
                        "reason": "Not evaluated by LLM — treated as unsatisfied.",
                    })

            # Enrich every evaluation with policy_source, policy_type, source_company from BQ
            for e in evaluations:
                p = policy_lookup.get(e.get("policy_id"), {})
                e["policy_source"] = p.get("policy_source", "")
                e["policy_type"] = p.get("policy_type", "regulatory")
                e["source_company"] = p.get("source_company") or None

            satisfied = [e for e in evaluations if e.get("satisfied")]
            unsatisfied = [e for e in evaluations if not e.get("satisfied")]
            score_pct = round(len(satisfied) / total * 100, 1) if total > 0 else 0.0

            print(f"[ROI] {industry}: {len(satisfied)}/{total} policies satisfied ({score_pct}%)", flush=True)

            return {
                "industry": industry,
                "action_evaluated": action_description,
                "total_policies": total,
                "policies_satisfied": len(satisfied),
                "eligibility_score_pct": score_pct,
                "satisfied_policies": satisfied,
                "unsatisfied_policies": unsatisfied,
            }

        except Exception as e:
            print(f"[ROI] Procurement ROI calculation failed: {e}", flush=True)
            return {"error": str(e)}

    def get_policy_gaps(self, user_id: str) -> dict:
        """
        Identifies which procurement policies for the user's industry are NOT yet covered
        by any of their existing sustainability actions. Returns a prioritised gap list,
        coverage percentage, and top 5 gaps to address first.
        """
        try:
            session_id = f"gaps_{datetime.now().strftime('%Y%m%d')}"
            fs = FireStoreChat(user_id, session_id)
            user_doc = fs.user_ref.get()
            if not user_doc.exists:
                return {"error": "User not found"}

            user_data = user_doc.to_dict()
            industry = user_data.get("company_industry", "").strip()
            if not industry:
                return {"error": "User's industry is not set in their profile"}

            dataset_id = "dash_beta_database"
            table_id = "procurement_policies"
            query = f"""
                SELECT policy_id, policy_name, policy_description, category, policy_source,
                       IFNULL(policy_type, 'regulatory') as policy_type,
                       source_company
                FROM `{self.project_id}.{dataset_id}.{table_id}`
                WHERE LOWER(industry) = @industry
                ORDER BY policy_type, category, policy_name
            """
            job_config = bigquery.QueryJobConfig(
                query_parameters=[
                    bigquery.ScalarQueryParameter("industry", "STRING", industry.lower())
                ]
            )
            results = list(self.bq_client.query(query, job_config=job_config).result())

            if not results:
                return {
                    "error": f"No procurement policies found for industry: {industry}",
                    "hint": "Run scripts/populate_procurement_policies.py to seed policies for this industry."
                }

            policies = [
                {
                    "policy_id": row.policy_id,
                    "policy_name": row.policy_name,
                    "policy_description": row.policy_description,
                    "category": row.category,
                    "policy_source": row.policy_source,
                    "policy_type": row.policy_type,
                    "source_company": row.source_company or None,
                }
                for row in results
            ]
            policy_lookup = {p["policy_id"]: p for p in policies}

            actions_result = api_client.view_actionList(user_id)
            existing_actions = actions_result.get("actions", [])

            if not existing_actions:
                print(f"[PolicyGaps] No existing actions for user {user_id}; all {len(policies)} policies are gaps.", flush=True)
                return {
                    "industry": industry,
                    "total_policies": len(policies),
                    "covered_count": 0,
                    "gap_count": len(policies),
                    "coverage_pct": 0.0,
                    "gaps": policies,
                    "covered": [],
                    "top_gaps": policies[:5],
                    "note": "No sustainability actions found. All policies are currently unaddressed — this is your starting point.",
                }

            action_summaries = [
                f"- {a.get('action_name', 'Unknown')}: {a.get('action_description', '')}"
                for a in existing_actions
            ]
            actions_text = "\n".join(action_summaries)

            prompt = f"""You are assessing which procurement policies are already covered by a company's existing sustainability actions.

Industry: {industry}

Existing sustainability actions:
{actions_text}

Procurement policies to assess:
{json.dumps(policies, indent=2)}

For each policy, determine whether at least one of the existing actions clearly and substantially addresses that policy.
Mark 'covered' as true only if an existing action clearly and substantially fulfils that policy requirement — partial alignment does not count.

You MUST return an evaluation for every policy in the list. Do not skip any.

Return ONLY a valid JSON object in this exact structure:
{{
  "evaluations": [
    {{
      "policy_id": "<id>",
      "policy_name": "<name>",
      "category": "<category>",
      "covered": true or false,
      "covering_action": "<name of the action that covers it, or null if not covered>",
      "reason": "<one sentence explanation>"
    }}
  ]
}}"""

            evaluations = []
            for attempt in range(3):
                try:
                    response = self.genai_client.models.generate_content(
                        model="gemini-3.5-flash",
                        contents=prompt,
                        config={"response_mime_type": "application/json"}
                    )
                    evaluations = json.loads(response.text).get("evaluations", [])
                    break
                except Exception as llm_err:
                    print(f"[PolicyGaps] LLM eval attempt {attempt + 1} failed: {llm_err}", flush=True)
                    if attempt == 2:
                        return {"error": f"LLM evaluation failed: {llm_err}"}
                    time.sleep(2 ** attempt)

            evaluated_ids = {e.get("policy_id") for e in evaluations}
            for p in policies:
                if p["policy_id"] not in evaluated_ids:
                    evaluations.append({
                        "policy_id": p["policy_id"],
                        "policy_name": p["policy_name"],
                        "category": p["category"],
                        "covered": False,
                        "covering_action": None,
                        "reason": "Not evaluated by LLM — treated as a gap.",
                    })

            for e in evaluations:
                pol = policy_lookup.get(e.get("policy_id"), {})
                e["policy_source"] = pol.get("policy_source", "")
                e["policy_description"] = pol.get("policy_description", "")

            covered = [e for e in evaluations if e.get("covered")]
            gaps = [e for e in evaluations if not e.get("covered")]
            total = len(policies)
            coverage_pct = round(len(covered) / total * 100, 1) if total > 0 else 0.0

            print(f"[PolicyGaps] {industry}: {len(covered)}/{total} policies covered ({coverage_pct}%)", flush=True)

            return {
                "industry": industry,
                "total_policies": total,
                "covered_count": len(covered),
                "gap_count": len(gaps),
                "coverage_pct": coverage_pct,
                "gaps": gaps,
                "covered": covered,
                "top_gaps": gaps[:5],
                "note": f"Found {len(gaps)} policy gaps. Suggestions should prioritise the top_gaps list.",
            }

        except Exception as e:
            print(f"[PolicyGaps] Failed: {e}", flush=True)
            return {"error": str(e)}

    def get_decarbonisation_actions(self, user_id: str) -> dict:
        """
        Retrieves real-world decarbonisation initiatives and actions from sustainability
        reports of leading companies in the user's industry. Use these as concrete, grounded
        ideas for proposing sustainability actions to the user to satisfy procurement policies.
        """
        try:
            session_id = f"decarb_{datetime.now().strftime('%Y%m%d')}"
            fs = FireStoreChat(user_id, session_id)
            user_doc = fs.user_ref.get()
            if not user_doc.exists:
                return {"error": "User not found"}

            user_data = user_doc.to_dict()
            industry = user_data.get("company_industry", "").strip()
            if not industry:
                return {"error": "User's industry is not set in their profile"}

            dataset_id = "dash_beta_database"
            table_id = "decarbonisation_plans"
            query = f"""
                SELECT plan_id, company_name, action_name, action_description, category, target_year, source_report
                FROM `{self.project_id}.{dataset_id}.{table_id}`
                WHERE LOWER(industry) = @industry
                ORDER BY category, company_name, action_name
            """
            job_config = bigquery.QueryJobConfig(
                query_parameters=[
                    bigquery.ScalarQueryParameter("industry", "STRING", industry.lower())
                ]
            )
            results = list(self.bq_client.query(query, job_config=job_config).result())

            actions = [
                {
                    "plan_id": row.plan_id,
                    "company_name": row.company_name,
                    "action_name": row.action_name,
                    "action_description": row.action_description,
                    "category": row.category,
                    "target_year": row.target_year,
                    "source_report": row.source_report,
                }
                for row in results
            ]

            print(f"[DecarbonisationActions] {industry}: Found {len(actions)} actions", flush=True)

            return {
                "industry": industry,
                "total_actions": len(actions),
                "decarbonisation_actions": actions,
                "note": f"Ground suggested sustainability actions for the user in these real-world initiatives from other {industry} companies.",
            }

        except Exception as e:
            print(f"[DecarbonisationActions] Failed: {e}", flush=True)
            return {"error": str(e)}

    def get_industry_guidelines(self, industry_name: str) -> dict:
        """
        Fetches industry-specific procurement policies and recommended actions from BigQuery.
        """
        try:
            dataset_id = "dash_beta_database"
            table_id = "industry_guidelines"
            query = f"""
                SELECT policy_document, predefined_actions, key_questions 
                FROM `{self.project_id}.{dataset_id}.{table_id}` 
                WHERE LOWER(industry_name) = @industry
            """
            job_config = bigquery.QueryJobConfig(
                query_parameters=[
                    bigquery.ScalarQueryParameter("industry", "STRING", industry_name.lower())
                ]
            )
            
            query_job = self.bq_client.query(query, job_config=job_config)
            results = list(query_job.result())

            if not results:
                return {"message": f"No specific guidelines found for industry: {industry_name}"}

            row = results[0]
            
            # Parse JSON fields if they are stored as strings
            def safe_json_load(data):
                if isinstance(data, str):
                    try:
                        import json
                        return json.loads(data)
                    except:
                        return data
                return data

            return {
                "industry": industry_name,
                "policy_document": row.policy_document,
                "predefined_actions": safe_json_load(row.predefined_actions),
                "key_questions": safe_json_load(row.key_questions)
            }

        except Exception as e:
            print(f"Error fetching industry guidelines: {e}")
            return {"error": str(e)}
        
    def add_action(
            self,
            user_id: str,
            action_id: str,
            action_name: str,
            action_type: str,
            action_description: str,
            estimated_spend: float,
            estimated_co2_reduced: float,
            estimated_revenue_unlocked: float,
            plan_id: str,
            timeline_start: datetime,
            timeline_end: datetime,
            status: str
    ): 
        """Adds a new sustainability action to the database. This function returns the numbers of rows added to the database. Use this when you need to add a new action to the database."""
        try:
            print("Adding new action")
            action_payload = {
                "action_id": action_id,
                "action_name": action_name,
                "action_type": action_type,
                "action_description": action_description,
                "estimated_spend": estimated_spend,
                "estimated_co2_reduced": estimated_co2_reduced,
                "estimated_revenue_unlocked": estimated_revenue_unlocked,
                "plan_id": plan_id,
                "timeline_start": timeline_start.isoformat(),
                "timeline_end": timeline_end.isoformat(),
                "status": status
            }
            return api_client.add_action_service(user_id, action_payload)
        except Exception as e:
            return {
                "error": {e}
            }
        
    def remove_action(
            self,
            user_id: str,
            action_id: str
    ):
        """Removes a sustainability action from the database. Use this when you need to delete an action from the database."""
        try:
            return api_client.remove_action_service(user_id, action_id)
        except Exception as e:
            return {
                "error": {e}
            }
        
    def update_action(
            self,
            user_id: str,
            action_id: str,
            actual_co2_reduced: float,
            actual_spend: float,
            actual_revenue_unlocked: float,
            day_started: datetime,
            day_completed: datetime,
    ):
        """Updates an existing sustainability action in the database. Use this when you need to update details of an existing action."""
        try:
            update_action_payload = {
                "co2_red": actual_co2_reduced,
                "spend": actual_spend,
                "rev_unlocked": actual_revenue_unlocked,
                "day_start": day_started.isoformat() if day_started else None,
                "day_end": day_completed.isoformat() if day_completed else None,
            }
            return api_client.update_action_service(user_id, action_id, update_action_payload)
        except Exception as e:
            return {
                "error": {e}
            }

    def vertex_search(self, query: str, user_id: str):
        """
        Search user's sustainability documents (bills, emissions data, etc.)
        to answer specific questions about their energy or environmental impact.

        Args:
            query: Natural language question (e.g. 'What was my electricity cost in March?')
            user_id: The ID of the user whose documents to search.
        """
        self._emit_status("vertex_search")
        # The full resource name of the search app serving config
        serving_config = f"projects/{self.project_id}/locations/{self.location_vertexAI}/collections/default_collection/engines/{self.engine_id}/servingConfigs/default_config"

        query_understanding_spec = discoveryengine.AnswerQueryRequest.QueryUnderstandingSpec(
            query_rephraser_spec=discoveryengine.AnswerQueryRequest.QueryUnderstandingSpec.QueryRephraserSpec(
                disable=False,  # Disable query rephraser
                max_rephrase_steps=1,  # Number of rephrase steps
            ),
            # Optional: Classify query types
            query_classification_spec=discoveryengine.AnswerQueryRequest.QueryUnderstandingSpec.QueryClassificationSpec(
                types=[
                    discoveryengine.AnswerQueryRequest.QueryUnderstandingSpec.QueryClassificationSpec.Type.ADVERSARIAL_QUERY,
                    discoveryengine.AnswerQueryRequest.QueryUnderstandingSpec.QueryClassificationSpec.Type.NON_ANSWER_SEEKING_QUERY,
                ]  # Options: ADVERSARIAL_QUERY, NON_ANSWER_SEEKING_QUERY or both
            ),
        )

        answer_generation_spec = discoveryengine.AnswerQueryRequest.AnswerGenerationSpec(
            ignore_adversarial_query=False,  # Ignore adversarial query
            ignore_non_answer_seeking_query=False,  # Ignore non-answer seeking query
            ignore_low_relevant_content=False,  # Return fallback answer when content is not relevant
            model_spec=discoveryengine.AnswerQueryRequest.AnswerGenerationSpec.ModelSpec(
                model_version="stable",  # Model to use for answer generation
            ),
            prompt_spec=discoveryengine.AnswerQueryRequest.AnswerGenerationSpec.PromptSpec(
                preamble="""    
                Given the conversation between a user and a helpful assistant and some search results, create a final answer for the assistant. 
                The answer should use all relevant information from the search results, not introduce any additional information, and use exactly the same words as the search results when possible.
                """, 
            ),
            include_citations=True,  # Include citations in the response
            answer_language_code="en", 
        )

        request = discoveryengine.AnswerQueryRequest(
            serving_config=serving_config,
            query=discoveryengine.Query(text=query),
            session=None,  # Include previous session ID to continue a conversation
            query_understanding_spec=query_understanding_spec,
            answer_generation_spec=answer_generation_spec,
            user_pseudo_id=user_id,  # Add user pseudo-identifier for queries.
        )

        # Make the request
        response = self.discovery_engine_client.answer_query(request)

        return self._format_vertex_response(response)

    def _format_vertex_response(self, response):
        """Clean up the Vertex AI Search response for the LLM."""
        formatted = {
            "answer": response.answer.answer_text if response.answer else "No answer found.",
            "citations": [],
            "source_documents": []
        }

        # Extract citations
        if response.answer and response.answer.citations:
            for citation in response.answer.citations:
                formatted["citations"].append({
                    "start": citation.start_index,
                    "end": citation.end_index,
                    "reference_indices": [s.reference_id for s in citation.sources]
                })

        # Extract unique source documents and their metadata
        seen_docs = set()
        for ref in response.answer.references:
            if not ref.chunk_info or not ref.chunk_info.document_metadata:
                continue
                
            doc_meta = ref.chunk_info.document_metadata
            if doc_meta.document not in seen_docs:
                seen_docs.add(doc_meta.document)
                
                # Convert struct_data fields to a simple dict
                metadata = {}
                if doc_meta.struct_data:
                    for key, value in doc_meta.struct_data.items():
                        # Handle the nested 'data' (JSON) field or any other sub-structs
                        if hasattr(value, "items"): # It's a MapComposite or dict
                            nested = {}
                            for n_key, n_val in value.items():
                                if hasattr(n_val, "string_value"):
                                    nested[n_key] = n_val.string_value
                                elif hasattr(n_val, "number_value"):
                                    nested[n_key] = n_val.number_value
                                else:
                                    nested[n_key] = n_val
                            metadata[key] = nested
                        elif hasattr(value, "string_value"):
                            metadata[key] = value.string_value
                        elif hasattr(value, "number_value"):
                            metadata[key] = value.number_value
                        else:
                            metadata[key] = value

                formatted["source_documents"].append({
                    "id": doc_meta.document.split("/")[-1],
                    "title": doc_meta.title,
                    "uri": doc_meta.uri,
                    "metadata": metadata
                })
                
        return formatted

    def fetch_octopus_usage(self, user_id: str, days_back: int = 7, period_from: str = None, period_to: str = None):
        """Fetches electricity usage and cost for a user from Octopus Energy. Supports specific date ranges."""
        self._emit_status("octopus_fetch")
        try:
            print(f"[Tool] Fetching energy data from Octopus for user: {user_id}, range: {period_from} to {period_to}")
            # 1. Fetch energy settings from Firestore
            session_id = f"{datetime.now().strftime('%Y%m%d')}"
            firestore = FireStoreChat(user_id, session_id) 
            settings = firestore.get_user_energy_context()
            
            if not settings:
                return f"No energy settings (account number) found for user {user_id} in Firestore."
            
            account_number = settings.get("account_number")
            secret_name = settings.get("octopus_secret_name")
            
            if not all([account_number, secret_name]):
                return f"Incomplete energy settings found for user {user_id}. Need account number and secret_name."

            # 2. Fetch from API
            client = OctopusClient(self.project_id)
            api_key = client.get_secret(secret_name)
            account_data = client.get_account_details(account_number=account_number, api_key=api_key)
            mpan_list = account_data["mpan_list"]
            start_date = account_data["start_date"]

            if period_from and period_from < start_date:
                return f"No data available for the selected period. Please select a period after {start_date}"
            if period_to and period_to > datetime.now().strftime("%Y-%m-%d"):
                period_to = datetime.now().strftime("%Y-%m-%d")

            energy_data_list = []

            for mpan in mpan_list:
                for serial in mpan_list[mpan]:
                    energy_data = client.get_summarized_usage(user_id, mpan, serial, secret_name, days_back, period_from, period_to)
                    if not energy_data:
                        print(f"No new data for user {user_id}.")
                        continue

                    energy_data_list.append(energy_data)
            
            if not energy_data_list:
                return "No energy consumption data found for the requested period. Ask user if they want to add new sources to cover for this period."
            
            # Format a summary for the agent to read
            summary = "Octopus Energy Report:\n"
            for data in energy_data_list:
                summary += (
                    f"- Meter {data.meter_serial} ({data.mpan}): "
                    f"{data.consumption_kwh:.2f} kWh from {data.period_start[:10]} to {data.period_end[:10]}\n"
                    f"Total Cost: £{data.total_cost_gbp:.2f}\n"
                )

            return summary

        except Exception as e:
            print(f"[Tool] Failed to get report from Octopus Energy: {str(e)}", flush=True)
            return f"Error fetching from Octopus: {str(e)}"

    def _generate_pdf_report(self, title: str, content: str, user_id: str, company_name: str, report_type: str) -> str:
        """
        Generates a PDF report, uploads it to GCS, and returns the URL.
        Filename format: companyName_reportType.pdf
        """
        try:
            print(f"[Tool] Generating PDF report for {company_name}: {report_type}", flush=True)
            
            # Custom PDF Class for ESG formatting
            class ESGReportPDF(FPDF):
                def header(self):
                    # Logo
                    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                    logo_path = os.path.join(base_dir, "chat_agent", "data", "logo.png")
                    try:
                        self.image(logo_path, 10, 8, 40)
                    except Exception as e:
                        print(f"Could not load logo: {e}")
                    
                    # Align vertically with logo and right-align the dynamic report type
                    self.set_y(12)
                    self.set_font('helvetica', 'B', 12)
                    self.set_text_color(5, 150, 105) # Dash Green #059669
                    self.cell(0, 10, report_type.upper(), 0, 0, 'R')
                    self.ln(25)

                def footer(self):
                    self.set_y(-15)
                    self.set_font('helvetica', 'I', 8)
                    self.set_text_color(128, 128, 128)
                    self.cell(0, 10, f'Page {self.page_no()}/{{nb}}', 0, 0, 'C')
                    
                    # Date alignment to the right
                    date_str = datetime.now().strftime("%B %d, %Y")
                    # Move X to right edge minus margin and text width
                    self.set_x(-60) 
                    self.cell(50, 10, f'Generated on {date_str}', 0, 0, 'R')

            pdf = ESGReportPDF()
            pdf.alias_nb_pages()
            pdf.add_page()
            
            # Report Title
            pdf.set_font("helvetica", "B", 18)
            pdf.set_text_color(33, 37, 41)
            pdf.cell(0, 10, title, new_x="LMARGIN", new_y="NEXT", align="L")
            pdf.ln(2)
            
            # Subtitle (Company & Report Type)
            pdf.set_font("helvetica", "I", 12)
            pdf.set_text_color(100, 100, 100)
            pdf.cell(0, 10, f"Prepared for: {company_name} | Type: {report_type}", new_x="LMARGIN", new_y="NEXT", align="L")
            pdf.ln(8)
            
            # Main Content
            pdf.set_font("helvetica", "", 11)
            pdf.set_text_color(50, 50, 50)
            clean_content = content.encode('latin-1', 'replace').decode('latin-1')
            pdf.multi_cell(0, 6, clean_content)
            
            # Sanitize and construct filename
            safe_company = company_name.replace(" ", "").replace("/", "_")
            safe_type = report_type.replace(" ", "").replace("/", "_")
            filename = f"{safe_company}_{safe_type}.pdf"
            
            temp_path = f"/tmp/{filename}"
            pdf.output(temp_path)
            
            # Upload to GCS
            bucket_name = f"{self.project_id}.firebasestorage.app"
            storage_client = storage.Client(project=self.project_id)
            bucket = storage_client.bucket(bucket_name)
            blob = bucket.blob(f"users/{user_id}/reports/{filename}")
            
            blob.upload_from_filename(temp_path)
            
            # Construct the public URL
            url = f"https://storage.googleapis.com/{bucket_name}/users/{user_id}/reports/{filename}"
            
            # Cleanup
            if os.path.exists(temp_path):
                os.remove(temp_path)
            
            print(f"[Tool] PDF report generated and uploaded: {filename}", flush=True)
            return f"PDF report '{filename}' generated successfully. You can access it in the Knowledge page."
            
            print(f"Error generating PDF: {e}", flush=True)
            return f"Failed to generate PDF. Please try again."

        except Exception as e:
            print(f"Error generating PDF: {e}", flush=True)
            return f"Failed to generate PDF. Please try again."

    def _pptx_set_bg(self, slide, colour) -> None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = colour

    def _pptx_add_textbox(self, slide, text: str, left, top, width, height,
                          font_size: int, bold: bool, colour,
                          align=None, word_wrap: bool = True) -> None:
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        if align is None:
            align = PP_ALIGN.LEFT
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = "Calibri"

    def _generate_pptx_report(
        self,
        title: str,
        subtitle: str,
        slides: list[dict],
        user_id: str,
        company_name: str,
        report_type: str,
    ) -> str:
        """
        Generates a branded PowerPoint presentation, uploads it to GCS, and returns
        a confirmation message. Each slide dict must have 'title', 'bullets' (list[str]),
        and optionally 'notes' (str).
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            DASH_GREEN      = RGBColor(0x05, 0x96, 0x69)
            DASH_GREEN_DARK = RGBColor(0x06, 0x4E, 0x3B)
            SLIDE_BG        = RGBColor(0xFF, 0xFF, 0xFF)
            BODY_TEXT       = RGBColor(0x1F, 0x29, 0x37)
            WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
            MINT            = RGBColor(0xA7, 0xF3, 0xD0)

            print(f"[Tool] Generating PPTX for {company_name}: {report_type}", flush=True)

            W = Inches(13.33)   # widescreen 16:9
            H = Inches(7.5)

            prs = Presentation()
            prs.slide_width  = W
            prs.slide_height = H

            blank_layout = prs.slide_layouts[6]  # blank

            # ── COVER SLIDE ──────────────────────────────────────────────
            cover = prs.slides.add_slide(blank_layout)
            self._pptx_set_bg(cover, DASH_GREEN_DARK)

            # Green accent bar on left
            bar = cover.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.35), H)
            bar.fill.solid()
            bar.fill.fore_color.rgb = DASH_GREEN
            bar.line.fill.background()

            # Logo (top-right, small)
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            logo_path = os.path.join(base_dir, "chat_agent", "data", "logo.png")
            if os.path.exists(logo_path):
                cover.shapes.add_picture(logo_path, Inches(11.8), Inches(0.25), width=Inches(1.2))

            # Title
            self._pptx_add_textbox(
                cover, title,
                left=Inches(0.7), top=Inches(2.4), width=Inches(11.0), height=Inches(1.6),
                font_size=40, bold=True, colour=WHITE, align=PP_ALIGN.LEFT,
            )

            # Subtitle
            if subtitle:
                self._pptx_add_textbox(
                    cover, subtitle,
                    left=Inches(0.7), top=Inches(4.1), width=Inches(11.0), height=Inches(0.8),
                    font_size=20, bold=False, colour=MINT, align=PP_ALIGN.LEFT,
                )

            # Company + date line
            date_str = datetime.now().strftime("%B %Y")
            self._pptx_add_textbox(
                cover, f"{company_name}  ·  {date_str}",
                left=Inches(0.7), top=Inches(6.6), width=Inches(10.0), height=Inches(0.5),
                font_size=13, bold=False, colour=MINT, align=PP_ALIGN.LEFT,
            )

            # ── CONTENT SLIDES ───────────────────────────────────────────
            for slide_data in slides:
                # slide_data may be a PPTXSlide Pydantic object or a plain dict
                if hasattr(slide_data, "title"):
                    s_title   = slide_data.title
                    s_bullets = slide_data.bullets
                    s_notes   = slide_data.notes or ""
                else:
                    s_title   = slide_data.get("title", "")
                    s_bullets = slide_data.get("bullets", [])
                    s_notes   = slide_data.get("notes", "")

                sl = prs.slides.add_slide(blank_layout)
                self._pptx_set_bg(sl, SLIDE_BG)

                # Green header band
                header_band = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.2))
                header_band.fill.solid()
                header_band.fill.fore_color.rgb = DASH_GREEN
                header_band.line.fill.background()

                # Slide title (white, on green band)
                self._pptx_add_textbox(
                    sl, s_title,
                    left=Inches(0.4), top=Inches(0.15), width=Inches(12.0), height=Inches(0.9),
                    font_size=24, bold=True, colour=WHITE, align=PP_ALIGN.LEFT,
                )

                # Bullet body
                if s_bullets:
                    body_box = sl.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(12.33), Inches(5.7))
                    tf = body_box.text_frame
                    tf.word_wrap = True

                    for i, bullet in enumerate(s_bullets):
                        if i == 0:
                            para = tf.paragraphs[0]
                        else:
                            para = tf.add_paragraph()
                        para.space_before = Pt(4)
                        run = para.add_run()
                        run.text = f"•  {bullet}"
                        run.font.size = Pt(18)
                        run.font.bold = False
                        run.font.color.rgb = BODY_TEXT
                        run.font.name = "Calibri"

                # Speaker notes
                if s_notes:
                    sl.notes_slide.notes_text_frame.text = s_notes

            # ── CLOSING SLIDE ────────────────────────────────────────────
            closing = prs.slides.add_slide(blank_layout)
            self._pptx_set_bg(closing, DASH_GREEN_DARK)
            bar2 = closing.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.35), H)
            bar2.fill.solid()
            bar2.fill.fore_color.rgb = DASH_GREEN
            bar2.line.fill.background()
            self._pptx_add_textbox(
                closing, "Thank you",
                left=Inches(0.7), top=Inches(2.8), width=Inches(11.0), height=Inches(1.0),
                font_size=36, bold=True, colour=WHITE, align=PP_ALIGN.LEFT,
            )
            self._pptx_add_textbox(
                closing, f"Prepared by Dash  ·  {company_name}",
                left=Inches(0.7), top=Inches(3.9), width=Inches(11.0), height=Inches(0.6),
                font_size=16, bold=False, colour=MINT, align=PP_ALIGN.LEFT,
            )

            # ── SAVE & UPLOAD ─────────────────────────────────────────────
            safe_company = company_name.replace(" ", "").replace("/", "_")
            safe_type    = report_type.replace(" ", "").replace("/", "_")
            filename     = f"{safe_company}_{safe_type}.pptx"
            temp_path    = f"/tmp/{filename}"
            prs.save(temp_path)

            bucket_name    = f"{self.project_id}.firebasestorage.app"
            storage_client = storage.Client(project=self.project_id)
            bucket         = storage_client.bucket(bucket_name)
            blob           = bucket.blob(f"users/{user_id}/reports/{filename}")
            blob.upload_from_filename(temp_path)

            if os.path.exists(temp_path):
                os.remove(temp_path)

            print(f"[Tool] PPTX generated and uploaded: {filename}", flush=True)
            return f"PowerPoint presentation '{filename}' generated successfully. You can access it in the Knowledge page."

        except Exception as e:
            print(f"[Tool] Error generating PPTX: {e}", flush=True)
            return f"Failed to generate PowerPoint. Please try again."

    def _calculate_emissions_batch_bq(self, activities: list[dict], user_id: str) -> dict:
        """
        Analytical calculation: Loads activities to a temp BQ table and joins with factors.
        This is the high-performance path for large datasets.
        """
        try:
            print(f"[Tool] Batch calculating emissions for {len(activities)} activities", flush=True)
            # Staging dataset for temp tables
            staging_dataset = "carbon_data"
            # Sanitize user_id for table name
            safe_user_id = user_id.replace("-", "_").replace(".", "_")
            temp_table_id = f"temp_staging_{safe_user_id}"
            full_table_name = f"{self.project_id}.{staging_dataset}.{temp_table_id}"
            
            # 1. Load data to BQ
            job_config = bigquery.LoadJobConfig(
                write_disposition="WRITE_TRUNCATE",
                autodetect=True,
            )
            load_job = self.bq_client.load_table_from_json(activities, full_table_name, job_config=job_config)
            load_job.result() # Wait for completion
            
            # 2. Run Analytical Join using Semantic Search with Exact Unit Filtering
            query = f"""
                SELECT * FROM (
                    SELECT 
                        query.row_id,
                        query.context as original_context,
                        query.activity_name as search_activity,
                        query.amount as original_amount,
                        query.unit as original_unit,
                        base.full_description as matched_factor,
                        base.factor as conversion_factor,
                        base.scope,
                        base.category as factor_category,
                        CAST(query.amount AS FLOAT64) * base.factor as co2_kg,
                        query.source_doc_name,
                        query.cell_ref,
                        'https://view.officeapps.live.com/op/view.aspx?src=https%3A%2F%2Fassets.publishing.service.gov.uk%2Fmedia%2F6846a4f55e92539572806125%2Fghg-conversion-factors-2025-full-set.xlsx&wdOrigin=BROWSELINK' as emission_factor_source,
                        ROW_NUMBER() OVER (PARTITION BY query.row_id ORDER BY distance) as rn
                    FROM VECTOR_SEARCH(
                        TABLE `carbon_data.emission_factors`,
                        'embedding',
                        (
                            SELECT * FROM ML.GENERATE_EMBEDDING(
                                MODEL `carbon_data.embedding_model`,
                                (SELECT activity_name || ' (' || unit || ')' as content, * FROM `{full_table_name}`),
                                STRUCT('RETRIEVAL_QUERY' AS task_type)
                            )
                        ),
                        'ml_generate_embedding_result',
                        top_k => 1000
                    )
                    WHERE LOWER(base.unit) IN (
                        LOWER(query.unit), 
                        'passenger.' || LOWER(query.unit),
                        'vehicle.' || LOWER(query.unit),
                        LOWER(query.unit) || '-vehicle'
                    )
                    AND base.category NOT LIKE 'WTT-%'
                ) WHERE rn = 1
            """
            query_job = self.bq_client.query(query)
            df = query_job.to_dataframe()
            print(f"[Tool] BQ result: {len(df)} matched rows from {len(activities)} input activities", flush=True)
            if not df.empty:
                print(f"[Tool] BQ columns: {df.columns.tolist()}", flush=True)
                print(f"[Tool] BQ sample:\n{df.head(2).to_string()}", flush=True)

            if df.empty:
                return {
                    "total_emissions_kgCO2e": 0,
                    "breakdown": [],
                    "audit_sample": [],
                    "status": "no_matches",
                    "processed_count": len(activities),
                    "message": (
                        "No emission factors matched any of the activities. "
                        "This usually means the units in the file don't match any known emission factor units. "
                        "Check that unit values (e.g. kWh, km, litres) are standard."
                    ),
                }

            # 3. Generate Audit Trail
            timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
            audit_filename = f"users/{user_id}/reports/audit_{timestamp}.xlsx"
            audit_bucket_name = "dash-beta-e61d0.firebasestorage.app"

            # Build a human-readable Excel file
            column_rename = {
                "row_id": "Row ID",
                "original_context": "Context",
                "search_activity": "Activity",
                "original_amount": "Amount",
                "original_unit": "Unit",
                "matched_factor": "Matched Emission Factor",
                "conversion_factor": "Conversion Factor (kgCO2e per unit)",
                "scope": "Scope",
                "factor_category": "Category",
                "co2_kg": "CO2e (kg)",
                "source_doc_name": "Source Document",
                "cell_ref": "Source Cell",
                "emission_factor_source": "Emission Factor Source",
            }
            audit_df = df.rename(columns=column_rename)

            excel_buffer = io.BytesIO()
            with pd.ExcelWriter(excel_buffer, engine="openpyxl") as writer:
                audit_df.to_excel(writer, index=False, sheet_name="Audit Trail")
                ws = writer.sheets["Audit Trail"]

                from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
                header_font = Font(bold=True, color="FFFFFF")
                header_fill = PatternFill("solid", fgColor="2D6A4F")
                center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
                thin_border = Border(
                    left=Side(style="thin"),
                    right=Side(style="thin"),
                    top=Side(style="thin"),
                    bottom=Side(style="thin"),
                )
                link_font = Font(color="0563C1", underline="single")

                link_col_names = {"Emission Factor Source"}
                header_row = [cell.value for cell in ws[1]]
                link_col_indices = {
                    col_idx for col_idx, name in enumerate(header_row, start=1)
                    if name in link_col_names
                }

                for col_idx, cell in enumerate(ws[1], start=1):
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = center_align
                    cell.border = thin_border
                    col_letter = cell.column_letter
                    max_len = len(str(cell.value)) if cell.value else 10
                    for data_cell in ws.iter_rows(min_row=2, min_col=col_idx, max_col=col_idx):
                        for dc in data_cell:
                            if dc.value is not None:
                                max_len = max(max_len, len(str(dc.value)))
                                if col_idx in link_col_indices and str(dc.value).startswith("http"):
                                    dc.hyperlink = dc.value
                                    dc.value = "View Source"
                                    dc.font = link_font
                            dc.border = thin_border
                    ws.column_dimensions[col_letter].width = min(max_len + 4, 50)

                ws.freeze_panes = "A2"

            excel_bytes = excel_buffer.getvalue()

            # Upload to GCS
            storage_client = storage.Client(project=self.project_id)
            bucket = storage_client.bucket(audit_bucket_name)
            if not bucket.exists():
                bucket = storage_client.create_bucket(audit_bucket_name)

            blob = bucket.blob(audit_filename)
            blob.upload_from_string(
                excel_bytes,
                content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
            audit_uri = f"gs://{audit_bucket_name}/{audit_filename}"
            
            # 4. Format results
            grand_total = df['co2_kg'].sum()
            summary = []
            
            # Group by scope, category, and unit
            grouped = df.groupby(['scope', 'factor_category', 'original_unit']).agg(
                total_co2=('co2_kg', 'sum'),
                total_amount=('original_amount', 'sum'),
                count=('row_id', 'count')
            ).reset_index()
            
            for _, row in grouped.iterrows():
                summary.append({
                    "scope": row['scope'],
                    "category": row['factor_category'],
                    "total_co2_kg": round(row['total_co2'], 2),
                    "total_amount": round(row['total_amount'], 2),
                    "unit": row['original_unit'],
                    "count": row['count']
                })
            
            # Sample for the LLM
            audit_sample = df.head(3).to_dict(orient="records")
            
            # Cleanup temp table
            self.bq_client.delete_table(full_table_name, not_found_ok=True)
            
            return {
                "total_emissions_kgCO2e": round(grand_total, 2),
                "breakdown": summary,
                "audit_sample": audit_sample,
                "status": "success",
                "processed_count": len(activities)
            }
        except Exception as e:
            print(f"Error in batch calculation: {e}", flush=True)
            try:
                self.bq_client.delete_table(full_table_name, not_found_ok=True)
            except: pass
            return {"error": str(e)}

    def _process_sheet(self, sheet_name: str, df: pd.DataFrame, source_doc_name: str = "") -> list[dict]:
        """
        Infers schema via Gemini and extracts activity rows for a single sheet.
        Designed to be called concurrently across sheets.
        """
        if df.empty:
            return []

        # Dynamic header promotion for dirty corporate sheets (e.g. title banners at the top)
        # data_start_row: 1-based Excel row number where data begins (row 1 = header in normal sheets)
        data_start_row = 2  # default: header at Excel row 1, data at row 2
        for idx in range(min(5, len(df))):
            row_vals = [str(x).strip().lower() for x in df.iloc[idx].values if pd.notna(x)]
            header_keywords = [
                "kwh", "amount", "total", "supplier", "month", "class", "unit", "distance",
                "km", "miles", "date", "passenger", "reason", "departure", "destination", "hotel", "postcode"
            ]
            matches = sum(1 for v in row_vals if any(k in v for k in header_keywords))
            current_has_unnamed = any("unnamed" in str(c).lower() or pd.isna(c) for c in df.columns)

            if current_has_unnamed and (matches >= 2 or (matches >= 1 and any("kwh" in v or "total" in v for v in row_vals))):
                new_headers = []
                for col_idx, val in enumerate(df.iloc[idx].values):
                    new_headers.append(str(val).strip() if pd.notna(val) else f"Unnamed: {col_idx}")
                df.columns = new_headers
                df = df.iloc[idx + 1:].reset_index(drop=True)
                # idx rows were title rows, promoted header is at Excel row idx+2, data starts at idx+3
                data_start_row = idx + 3
                print(f"[Tool] Promoted headers for {sheet_name} at row {idx}: {new_headers}", flush=True)
                break

        print(f"[Tool] Inferring schema for sheet: {sheet_name}", flush=True)

        csv_sample = df.head(2).to_csv(index=False)
        prompt = f"""
        Analyze the following CSV sample from a spreadsheet sheet named "{sheet_name}".
        Identify how to calculate carbon emissions from this data.

        Return a JSON object with:
        - amount_column: The exact name of the column containing the numerical value to calculate emissions on (e.g. "Kilometers", "Liters", "Cost", "kWh", "Nights").
        - unit: The implied or explicit unit for that amount (e.g. "km", "litres", "gbp", "kWh", "nights"). If not obvious, infer it.
        - activity_columns: A list of exact column names that help identify the specific type of transport or route (e.g., ["Vehicle Type", "Class", "Fuel", "Departure", "Destination"]). CRITICAL: DO NOT include columns containing employee names, dates, passenger names, reasons for travel, or project codes! Locations ARE helpful to determine if a flight/train is domestic or international.

        CSV Sample:
        {csv_sample}

        JSON:
        """

        # Retry logic to handle Gemini API rate limits (e.g. 429 quota exceeded)
        schema = None
        for attempt in range(3):
            try:
                response = self.genai_client.models.generate_content(
                    model="gemini-3.5-flash",
                    contents=prompt,
                    config={"response_mime_type": "application/json"}
                )
                print(f"[Tool] Schema for sheet '{sheet_name}' generated on attempt {attempt+1}", flush=True)
                schema = json.loads(response.text)
                break
            except Exception as llm_err:
                print(f"[Tool] Error inferring schema for '{sheet_name}' (attempt {attempt+1}): {llm_err}", flush=True)
                if attempt == 2:
                    return []
                time.sleep(2 ** attempt)

        amount_col = schema.get("amount_column")
        unit_val = schema.get("unit", "unknown")
        act_cols = schema.get("activity_columns", [])

        # Compute virtual 'nights' column for hotel sheets with check-in/check-out dates
        checkin_col = next((c for c in df.columns if any(k in str(c).lower() for k in ["check in", "check-in", "checkin"])), None)
        checkout_col = next((c for c in df.columns if any(k in str(c).lower() for k in ["check out", "check-out", "checkout"])), None)

        if checkin_col and checkout_col:
            try:
                df[checkin_col] = pd.to_datetime(df[checkin_col], errors='coerce')
                df[checkout_col] = pd.to_datetime(df[checkout_col], errors='coerce')
                df = df.dropna(subset=[checkin_col, checkout_col])
                df['nights'] = (df[checkout_col] - df[checkin_col]).dt.days
                amount_col = 'nights'
                unit_val = 'room per night'
                hotel_col = next((c for c in df.columns if "hotel" in str(c).lower()), None)
                dest_col = next((c for c in df.columns if "destination" in str(c).lower()), None)
                act_cols = [c for c in [hotel_col, dest_col] if c] or [df.columns[0]]
            except Exception as e:
                print(f"[Tool] Failed to calculate virtual nights column for '{sheet_name}': {e}", flush=True)

        if not amount_col or amount_col not in df.columns:
            print(f"[Tool] Could not find amount column for '{sheet_name}'. Skipping.", flush=True)
            return []

        # Resolve Excel column letter for the amount column so we can write cell refs to the audit trail.
        # 'nights' is a virtual computed column with no real cell, so we fall back to "N/A".
        from openpyxl.utils import get_column_letter
        if amount_col == 'nights':
            amount_col_letter = None  # computed column, no source cell
        else:
            try:
                amount_col_letter = get_column_letter(df.columns.tolist().index(amount_col) + 1)
            except ValueError:
                amount_col_letter = None

        activities = []
        columns = df.columns.tolist()

        # Optimize performance by bypassing df.iterrows()
        for df_idx, row in enumerate(df.to_dict('records')):
            try:
                # Skip summary/total rows
                is_total_row = any(
                    str(row[col]).strip().lower() in {"total", "totals", "grand total", "subtotal", "sum", "total / average"}
                    or str(row[col]).strip().lower().startswith("total ")
                    or str(row[col]).strip().lower().endswith(" total")
                    for col in columns if pd.notna(row[col])
                )
                if is_total_row:
                    continue

                amt = float(row[amount_col])
                if pd.isna(amt) or amt <= 0:
                    continue

                desc_parts = [str(row[c]) for c in act_cols if c in columns and pd.notna(row[c])]
                activity_name = f"{sheet_name} - " + " - ".join(desc_parts) if desc_parts else sheet_name

                if any(t in sheet_name.lower() or t in activity_name.lower() for t in ["flight", "plane", "aviation", "air"]):
                    if not any(c in activity_name.lower() for c in ["economy", "business", "first", "premium", "average"]):
                        activity_name += " - Average passenger"

                row_context = " | ".join(f"{k}: {v}" for k, v in row.items() if pd.notna(v))

                excel_row = df_idx + data_start_row
                cell_ref = f"{amount_col_letter}{excel_row}" if amount_col_letter else f"Row {excel_row} (computed)"

                activities.append({
                    "row_id": f"{sheet_name}_{len(activities)}",
                    "activity_name": activity_name,
                    "amount": amt,
                    "unit": unit_val,
                    "context": row_context,
                    "source_doc_name": source_doc_name,
                    "cell_ref": cell_ref,
                })
            except Exception:
                continue

        print(f"[Tool] Sheet '{sheet_name}': extracted {len(activities)} activities", flush=True)
        return activities

    def _calculate_emissions_from_file_bq(self, user_id: str) -> dict:
        """
        Reads all Excel/CSV files from the user's GCS uploads folder and processes
        them via the BigQuery analytical engine. All files are processed together.
        """
        self._emit_status("carbon_file_read")
        try:
            bucket_name = f"{self.project_id}.firebasestorage.app"
            prefix = f"users/{user_id}/uploads/"
            storage_client = storage.Client(project=self.project_id)
            blobs = [
                b for b in storage_client.bucket(bucket_name).list_blobs(prefix=prefix)
                if not b.name.endswith("/") and b.name.lower().endswith((".xlsx", ".csv"))
            ]
            if not blobs:
                return {"error": "No Excel or CSV file found in your uploads folder. Please upload a file first before running the calculation."}

            print(f"[Tool] Found {len(blobs)} file(s) to process: {[b.name.split('/')[-1] for b in blobs]}", flush=True)

            # Build a flat dict of {label: DataFrame} across all files and their sheets
            all_sheets: dict[str, pd.DataFrame] = {}
            for blob in blobs:
                filename = blob.name.split("/")[-1]
                gcs_uri = f"gs://{bucket_name}/{blob.name}"
                print(f"[Tool] Reading: {gcs_uri}", flush=True)
                try:
                    content = self._download_from_gcs(gcs_uri)
                    if gcs_uri.lower().endswith(".xlsx"):
                        file_sheets = pd.read_excel(io.BytesIO(content), sheet_name=None)
                    else:
                        file_sheets = {"Data": pd.read_csv(
                            io.StringIO(content) if isinstance(content, str) else io.BytesIO(content)
                        )}
                    print(f"[Tool] Blob metadata for '{filename}': {blob.metadata}", flush=True)
                    original_name = (blob.metadata or {}).get("originalFileName") or filename
                    for sheet_name, df in file_sheets.items():
                        if not df.empty:
                            label = f"{filename} — {sheet_name}" if len(file_sheets) > 1 else filename
                            display_name = f"{original_name} — {sheet_name}" if len(file_sheets) > 1 else original_name
                            all_sheets[label] = (df, display_name)
                except Exception as e:
                    print(f"[Tool] Failed to read '{filename}': {e}", flush=True)

            if not all_sheets:
                return {"error": "All uploaded files are empty or could not be read."}

            print(f"[Tool] Processing {len(all_sheets)} sheet(s) across all files in parallel", flush=True)

            # Infer schema and extract activities for all sheets concurrently
            all_activities = []
            max_workers = min(len(all_sheets), 5)
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = {
                    executor.submit(self._process_sheet, label, df, name): label
                    for label, (df, name) in all_sheets.items()
                }
                for future in as_completed(futures):
                    label = futures[future]
                    try:
                        activities = future.result()
                        all_activities.extend(activities)
                    except Exception as e:
                        print(f"[Tool] Sheet '{label}' failed: {e}", flush=True)

            if not all_activities:
                return {"error": "Could not extract any valid activities from the uploaded files. Please check the file format."}

            print(f"[Tool] Consolidated {len(all_activities)} activities across all files. Sending to BigQuery engine.", flush=True)
            self._emit_status("carbon_calculation")
            return self._calculate_emissions_batch_bq(all_activities, user_id)

        except Exception as e:
            print(f"Error processing file for BQ: {e}", flush=True)
            return {"error": str(e)}


    def _check_bulk_readiness(self, user_id: str, expected_categories: list[str] = []) -> dict:
        """
        Scans a GCS folder to check if all necessary documents are available before bulk processing.
        """
        try:
            bucket_name = f"{self.project_id}.firebasestorage.app"
            prefix = f"users/{user_id}/uploads/"
            gcs_uri = f"gs://{bucket_name}/{prefix}"
            print(f"[Tool] Checking bulk readiness for: {gcs_uri}", flush=True)

            storage_client = storage.Client(project=self.project_id)
            bucket = storage_client.bucket(bucket_name)
            # List blobs with the given prefix
            blobs = list(bucket.list_blobs(prefix=prefix))
            
            if not blobs:
                return {"error": "No files found in the specified GCS path."}
            
            file_summary = []
            categories_found = set()
            
            for blob in blobs:
                # Basic check to skip folders (some GCS providers create empty objects ending in /)
                if blob.name.endswith("/"): continue 
                
                filename = blob.name.split("/")[-1]
                if not filename: continue # Skip if it's just the prefix
                
                size_kb = blob.size / 1024
                
                # Simple keyword matching for categories from filename
                category = "Other"
                keywords = {
                    "Electricity": ["elec", "bill", "utility", "power", "enel", "octopus"],
                    "Fuel": ["fuel", "petrol", "diesel", "gasoline", "gas station"],
                    "Travel": ["flight", "hotel", "travel", "train", "airline", "booking"],
                    "Water": ["water", "sewage"],
                    "Waste": ["waste", "recycling", "trash"]
                }
                
                # Dynamically register expected categories as keywords
                for exp_cat in expected_categories:
                    if exp_cat not in keywords:
                        keywords[exp_cat] = [exp_cat.lower().strip()]
                
                for cat, keys in keywords.items():
                    if any(k in filename.lower() for k in keys):
                        category = cat
                        categories_found.add(cat)
                        break
                
                file_summary.append({
                    "filename": filename,
                    "gcs_uri": f"gs://{bucket_name}/{blob.name}",
                    "category_guess": category,
                    "size_kb": round(size_kb, 2),
                    "last_modified": blob.updated.strftime("%Y-%m-%d")
                })
            
            missing_categories = [cat for cat in expected_categories if cat.lower() not in [c.lower() for c in categories_found]]
            
            return {
                "total_files": len(file_summary),
                "categories_detected": list(categories_found),
                "missing_expected_categories": missing_categories,
                "file_list": file_summary[:20], # Return a snippet for the agent
                "ready_to_process": len(missing_categories) == 0,
                "message": "Folder check complete."
            }

        except Exception as e:
            print(f"Error checking bulk readiness: {e}", flush=True)
            return {"error": str(e)}


    def get_tools(self) -> list[Tool]:
        google_search_tool = Tool(
            name="google_search",
            description="Search Google for current events and real-time facts.",
            args_schema=google_search_input,
            func=self._logged_search,
        )

        read_actions_tool = Tool(
            name = "read_actions",
            description="Reads the existing sustainability actions from the database",
            args_schema = get_api,
            func=self.read_actions
        )

        add_action_tool = StructuredTool(
            name="add_action",
            description="Adds a new sustainability action to the database. Call directly with required parameters. Do NOT use print() or wrap this call.",
            args_schema=AddActionInput,
            func=self.add_action
        )

        remove_action_tool = StructuredTool(
            name="remove_action",
            description="Removes a sustainability action from the database",
            args_schema=RemoveActionInput,
            func=self.remove_action
        )

        update_action_tool = StructuredTool(
            name="update_action",
            description="Updates an existing sustainability action in the database. Call directly with required parameters. Do NOT use print() or wrap this call.",
            args_schema=UpdateActionInput,
            func=self.update_action
        )
        
        document_read_tool = StructuredTool(
            name="document_read",
            description="""
            Extract text from files (PDF, TXT, CSV, XLSX).
            Use this to read specific sustainability reports, utility bills (PDF), or Octopus API data logs (TXT/CSV).
            For user-uploaded files, ALWAYS use the exact gcs_uri returned by check_bulk_readiness — never construct or guess a GCS URL yourself.
            Accepts: Public URLs, Firebase URLs, GCS Console URLs (storage.cloud.google.com), or gs:// URIs.
            """,
            args_schema=document_read_input,
            func=lambda document_url: self._document_read(document_url),
        )

        vertex_search_tool = StructuredTool(
            name="vertex_search",
            description=self.vertex_search.__doc__,
            args_schema=vertex_search_input,
            func=self.vertex_search
        )
        
        octopus_fetch_tool = StructuredTool(
            name="fetch_octopus_usage",
            description="""
            Fetch LIVE energy consumption data and cost directly from Octopus Energy API. 
            Use this if the user asks for energy data that isnt available in the database or very recent data (e.g. today or yesterday) or specifically asks for a fresh catch.
            """,
            args_schema=EnergyFetchInput,
            func=self.fetch_octopus_usage
        )

        calculate_roi_tool = StructuredTool(
            name="calculate_procurement_roi",
            description="""
            Calculates a procurement eligibility score for a sustainability action by checking how many
            publicly available industry procurement policies it satisfies. Returns a percentage score,
            a list of satisfied policies, and a list of unsatisfied policies.
            Use this when the user asks about the ROI or business value of a sustainability action,
            or before adding an action to show how it improves procurement eligibility.
            """,
            args_schema=ProcurementROIInput,
            func=self.calculate_procurement_roi
        )

        get_policy_gaps_tool = StructuredTool(
            name="get_policy_gaps",
            description="""
            Identifies which industry procurement policies are NOT yet covered by the user's existing
            sustainability actions. Returns a gap list, coverage percentage, and top 5 priority gaps.
            Call this at the start of every planning session to ground suggestions in real, verifiable
            policy gaps. Works for both new users (no actions yet — returns all policies as gaps to
            build a plan from scratch) and existing users (evaluates current actions against all policies
            and shows what is still missing).
            """,
            args_schema=PolicyGapsInput,
            func=self.get_policy_gaps
        )

        get_decarbonisation_actions_tool = StructuredTool(
            name="get_decarbonisation_actions",
            description="""
            Retrieves real-world decarbonisation initiatives and actions from sustainability reports
            of leading companies in the user's industry. Use these as concrete, grounded ideas for
            proposing sustainability actions to the user to satisfy procurement policies.
            """,
            args_schema=DecarbonisationActionsInput,
            func=self.get_decarbonisation_actions
        )

        # industry_guidelines_tool = StructuredTool(
        #     name="get_industry_guidelines",
        #     description="""
        #     Fetches industry-specific procurement policies and recommended actions.
        #     Use this tool to understand the constraints and standard actions for the user's specific industry.
        #     """,
        #     args_schema=IndustryInfoInput,
        #     func=self.get_industry_guidelines
        # )

        generate_pdf_tool = StructuredTool(
            name="generate_pdf_report",
            description="""
            Generates a professional PDF report from text content.
            Use this when the user specifically asks for a PDF version of a summary, report, or analysis.
            """,
            args_schema=PDFGeneratorInput,
            func=self._generate_pdf_report
        )

        generate_pptx_tool = StructuredTool(
            name="generate_pptx_presentation",
            description="""
            Generates a branded PowerPoint (.pptx) presentation and uploads it to the user's Knowledge store.
            Use this when the user asks for a presentation, slide deck, PowerPoint, or slides version of a plan,
            report, or analysis. Structure the content into 4-10 slides — each slide gets a short title and
            3-6 bullet points. Provide a cover title, optional subtitle, and a list of slide objects.
            """,
            args_schema=PPTXGeneratorInput,
            func=self._generate_pptx_report
        )

        check_readiness_tool = StructuredTool(
            name="check_bulk_readiness",
            description="""
            Checks a GCS folder for documents before starting a bulk sustainability analysis. 
            Identifies file types and detects categories (Electricity, Fuel, Travel, etc.) from filenames.
            Use this to confirm with the user that all data is present before starting a long process.
            """,
            args_schema=BulkReadinessInput,
            func=self._check_bulk_readiness
        )


        calculate_bulk_file_tool = StructuredTool(
            name="calculate_emissions_from_structured_file",
            description="""
            Calculates carbon emissions for a large structured file (CSV/Excel) in GCS.
            The file MUST have columns for:
            1. 'Activity' (or 'Item', 'Description'): e.g., 'Electricity', 'Petrol'
            2. 'Amount' (or 'Qty', 'Usage'): e.g., 500
            3. 'Unit' (or 'Measure', 'UOM'): e.g., 'kWh', 'litres'
            This uses BigQuery's analytical engine for high-performance processing of thousands of rows.
            """,
            args_schema=BulkProcessInput,
            func=self._calculate_emissions_from_file_bq
        )

        tools_list = [
            google_search_tool,
            document_read_tool,
            read_actions_tool,
            add_action_tool,
            remove_action_tool,
            update_action_tool,
            vertex_search_tool,
            octopus_fetch_tool,
            calculate_roi_tool,
            get_policy_gaps_tool,
            get_decarbonisation_actions_tool,
            # industry_guidelines_tool,
            generate_pdf_tool,
            generate_pptx_tool,
            check_readiness_tool,
            calculate_bulk_file_tool
        ]

        return tools_list

# if __name__ == '__main__':
#     tool = ToolList()
#     tool.vertex_search("Tell me about my utility bills", "CORZZX0MxTQtGyAD7PSCI1HLp3y2")